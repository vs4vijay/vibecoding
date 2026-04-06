"""PowerPoint reader module for loading PPTX files and extracting slide content."""
import os
import tempfile
import subprocess
from pathlib import Path
from typing import Optional

from PIL import Image

from pptx import Presentation

from slidespeak.core.models import SlideSettings


class PowerPointReader:
    """Reads PowerPoint files and extracts slide content with notes.
    
    Uses fallback chain: COM export (Windows) → LibreOffice → python-pptx
    """
    
    def __init__(self, pptx_path: str):
        self.pptx_path = pptx_path
        self.prs = Presentation(pptx_path)
    
    def get_slide_count(self) -> int:
        """Return the total number of slides."""
        return len(self.prs.slides)
    
    def extract_slides_and_notes(self) -> list[SlideSettings]:
        """Extract slide data and notes, returning list of SlideSettings.
        
        Returns:
            List of SlideSettings objects with index and notes populated.
        """
        slides_data = []
        
        for i, slide in enumerate(self.prs.slides):
            slide_data = SlideSettings(
                slide_index=i,
                notes=""
            )
            
            # Extract notes
            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                slide_data.notes = notes_frame.text.strip()
            
            slides_data.append(slide_data)
        
        return slides_data
    
    def get_slide_thumbnails(self, dpi: int = 150, size: Optional[tuple[int, int]] = None) -> list[Image.Image]:
        """Generate thumbnails for all slides.
        
        Args:
            dpi: Resolution for rendering (default 150)
            size: Optional (width, height) to resize thumbnails
            
        Returns:
            List of PIL Image objects representing slide thumbnails
        """
        images = self._export_slides_as_images(dpi)
        
        if size:
            images = [img.resize(size, Image.Resampling.LANCZOS) for img in images]
        
        return images
    
    def get_slide_as_image(self, slide_index: int, dpi: int = 150) -> Image.Image:
        """Get a single slide as a PIL Image.
        
        Args:
            slide_index: Index of the slide (0-based)
            dpi: Resolution for rendering (default 150)
            
        Returns:
            PIL Image of the specified slide
            
        Raises:
            IndexError: If slide_index is out of range
        """
        if slide_index < 0 or slide_index >= len(self.prs.slides):
            raise IndexError(f"Slide index {slide_index} out of range (0-{len(self.prs.slides) - 1})")
        
        images = self._export_slides_as_images(dpi)
        return images[slide_index]
    
    def _export_slides_as_images(self, dpi: int = 150) -> list[Image.Image]:
        """Export slides as PIL Images using fallback chain.
        
        Tries: COM export (Windows) → LibreOffice → python-pptx
        
        Args:
            dpi: Resolution for rendering
            
        Returns:
            List of PIL Image objects
        """
        with tempfile.TemporaryDirectory() as temp_dir:
            # Try COM export first (most reliable on Windows)
            try:
                return self._com_export(temp_dir, dpi)
            except Exception as e:
                print(f"COM export issue: {e}, trying LibreOffice...")
            
            # Try LibreOffice as backup
            try:
                return self._libreoffice_export(temp_dir, dpi)
            except Exception as e2:
                print(f"LibreOffice failed: {e2}")
            
            # Final fallback: use python-pptx
            return self._python_pptx_fallback(dpi)
    
    def _com_export(self, temp_dir: str, dpi: int) -> list[Image.Image]:
        """Export slides using Windows COM automation."""
        import comtypes.client
        import time
        
        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        powerpoint.Visible = 1
        
        abs_path = os.path.abspath(self.pptx_path)
        deck = powerpoint.Presentations.Open(abs_path)
        
        slide_count = deck.Slides.Count
        image_paths = []
        
        for i in range(1, slide_count + 1):
            slide = deck.Slides(i)
            output_path = os.path.join(temp_dir, f"slide_{i-1:03d}.png")
            slide.Export(output_path, "PNG")
            image_paths.append(output_path)
        
        deck.Close()
        powerpoint.Quit()
        
        # Verify we got actual images
        if image_paths and os.path.exists(image_paths[0]):
            img = Image.open(image_paths[0])
            pixels = list(img.getdata())
            white_count = sum(1 for p in pixels[:1000] if p[:3] == (255, 255, 255))
            if white_count > 900:  # If more than 90% white, COM failed
                raise Exception("COM export produced blank images")
        
        return [Image.open(p) for p in image_paths]
    
    def _libreoffice_export(self, temp_dir: str, dpi: int) -> list[Image.Image]:
        """Export slides via LibreOffice PDF conversion."""
        # Convert to PDF first
        pdf_path = os.path.join(temp_dir, "slides.pdf")
        result = subprocess.run([
            "soffice", "--headless", "--convert-to", "pdf",
            "--outdir", temp_dir, self.pptx_path
        ], capture_output=True, text=True, timeout=60)
        
        # Find the PDF
        base_name = os.path.splitext(os.path.basename(self.pptx_path))[0] + ".pdf"
        pdf_path = os.path.join(temp_dir, base_name)
        
        if not os.path.exists(pdf_path):
            raise Exception(f"LibreOffice did not produce PDF at {pdf_path}")
        
        from pdf2image import convert_from_path
        images = convert_from_path(pdf_path, dpi=dpi)
        return images
    
    def _python_pptx_fallback(self, dpi: int) -> list[Image.Image]:
        """Fallback export using python-pptx to render slides as images."""
        from PIL import Image, ImageDraw, ImageFont
        
        # Get slide dimensions (in inches)
        width_in = self.prs.slide_width / 914400
        height_in = self.prs.slide_height / 914400
        
        # Convert to pixels at specified DPI
        width_px = int(width_in * dpi)
        height_px = int(height_in * dpi)
        
        images = []
        
        for i, slide in enumerate(self.prs.slides):
            # Create image with light gray background
            img = Image.new('RGB', (width_px, height_px), '#f0f0f0')
            draw = ImageDraw.Draw(img)
            
            # Try to load a default font
            try:
                font = ImageFont.truetype("arial.ttf", 24)
                font_small = ImageFont.truetype("arial.ttf", 18)
            except:
                font = ImageFont.load_default()
                font_small = ImageFont.load_default()
            
            # Extract and draw title
            title_y = 50
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    try:
                        # Get position
                        left = int(shape.left / 914400 * dpi) if hasattr(shape, "left") else 50
                        top = int(shape.top / 914400 * dpi) if hasattr(shape, "top") else title_y
                        width = int(shape.width / 914400 * dpi) if hasattr(shape, "width") else width_px - 100
                        
                        # Determine font size based on shape position
                        current_font = font_small if top > 200 else font
                        
                        # Draw text
                        draw.text((left, top), shape.text[:100], fill='black', font=current_font)
                        title_y = top + 50
                    except Exception:
                        pass
            
            images.append(img)
        
        return images
