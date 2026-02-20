from pptx import Presentation

prs = Presentation('sample.pptx')

notes_texts = [
    'Welcome to this presentation about our business plan. Today we will discuss our market analysis, financial projections, and growth strategy.',
    'Our market research shows significant opportunity in the target segment. We have identified key customer demographics and their specific needs.',
    'Our product addresses these needs with innovative features and competitive pricing. We have already received positive feedback from early adopters.',
    'Our financial projections show strong growth over the next five years. We expect to break even by year two and achieve profitability by year three.',
    'Thank you for your attention. We are happy to answer any questions you may have.'
]

for i, slide in enumerate(prs.slides):
    if i < len(notes_texts):
        notes_slide = slide.notes_slide
        notes_frame = notes_slide.notes_text_frame
        notes_frame.text = notes_texts[i]

prs.save('sample_with_notes.pptx')
print('Added speaker notes to presentation')
