import os
import tempfile
import shutil
from pathlib import Path
from typing import Optional

import click
from pptx import Presentation
from gtts import gTTS
from moviepy import ImageClip, AudioFileClip, concatenate_videoclips
from moviepy.video.io.VideoFileClip import VideoFileClip
from PIL import Image


class PowerPointReader:
    def __init__(self, pptx_path: str):
        self.pptx_path = pptx_path
        self.prs = Presentation(pptx_path)

    def get_slide_count(self) -> int:
        return len(self.prs.slides)

    def extract_slides_and_notes(self) -> list[dict]:
        slides_data = []
        
        for i, slide in enumerate(self.prs.slides):
            slide_data = {
                "index": i,
                "notes": "",
                "has_content": False
            }
            
            # Extract notes
            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                slide_data["notes"] = notes_frame.text.strip()
            
            # Check if slide has any content (shapes, pictures, etc.)
            if len(slide.shapes) > 0:
                slide_data["has_content"] = True
            
            slides_data.append(slide_data)
        
        return slides_data


class TextToSpeech:
    def __init__(self, output_dir: str, lang: str = "en"):
        self.output_dir = output_dir
        self.lang = lang

    def text_to_audio(self, text: str, output_filename: str) -> str:
        if not text.strip():
            return None
            
        output_path = os.path.join(self.output_dir, output_filename)
        tts = gTTS(text=text, lang=self.lang)
        tts.save(output_path)
        return output_path


class SlideToImage:
    def __init__(self, pptx_path: str, output_dir: str, dpi: int = 150):
        self.pptx_path = pptx_path
        self.output_dir = output_dir
        self.dpi = dpi

    def export_slides(self) -> list[str]:
        image_paths = []
        
        # Try COM export first (most reliable on Windows)
        try:
            import comtypes.client
            import os
            import time
            
            powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
            powerpoint.Visible = 1
            
            abs_path = os.path.abspath(self.pptx_path)
            deck = powerpoint.Presentations.Open(abs_path)
            
            slide_count = deck.Slides.Count
            print(f"Exporting {slide_count} slides via PowerPoint COM...")
            
            for i in range(1, slide_count + 1):
                slide = deck.Slides(i)
                output_path = os.path.join(self.output_dir, f"slide_{i-1:03d}.png")
                
                slide.Export(output_path, "PNG")
                image_paths.append(output_path)
                print(f"  Exported slide {i}/{slide_count}")
            
            deck.Close()
            powerpoint.Quit()
            
            # Verify we got actual images
            if image_paths and os.path.exists(image_paths[0]):
                img = Image.open(image_paths[0])
                pixels = list(img.getdata())
                white_count = sum(1 for p in pixels[:1000] if p[:3] == (255, 255, 255))
                if white_count > 900:  # If more than 90% white, COM failed
                    print("COM export produced blank images, trying fallback...")
                    image_paths = []
                    raise Exception("COM export produced blank images")
            
            return image_paths
            
        except Exception as e:
            print(f"COM export issue: {e}, trying LibreOffice...")
            
            # Try LibreOffice as backup
            try:
                import subprocess
                import os
                
                # Convert to PDF first
                pdf_path = os.path.join(self.output_dir, "slides.pdf")
                result = subprocess.run([
                    "soffice", "--headless", "--convert-to", "pdf", 
                    "--outdir", self.output_dir, self.pptx_path
                ], capture_output=True, text=True, timeout=60)
                
                print(f"LibreOffice result: {result.returncode}")
                print(f"Stdout: {result.stdout}")
                print(f"Stderr: {result.stderr}")
                
                # Find the PDF
                base_name = os.path.splitext(os.path.basename(self.pptx_path))[0] + ".pdf"
                pdf_path = os.path.join(self.output_dir, base_name)
                
                if os.path.exists(pdf_path):
                    from pdf2image import convert_from_path
                    images = convert_from_path(pdf_path, dpi=self.dpi)
                    for i, img in enumerate(images):
                        output_path = os.path.join(self.output_dir, f"slide_{i:03d}.png")
                        img.save(output_path, "PNG")
                        image_paths.append(output_path)
                    return image_paths
                    
            except Exception as e2:
                print(f"LibreOffice failed: {e2}")
            
            # Final fallback: use python-pptx with better text extraction
            return self._better_fallback_export()
    
    def _better_fallback_export(self) -> list[str]:
        from pptx import Presentation
        from PIL import Image, ImageDraw, ImageFont
        
        image_paths = []
        prs = Presentation(self.pptx_path)
        
        # Get slide dimensions (in inches)
        width_in = prs.slide_width / 914400
        height_in = prs.slide_height / 914400
        
        # Convert to pixels at specified DPI
        width_px = int(width_in * self.dpi)
        height_px = int(height_in * self.dpi)
        
        print(f"Using fallback export: {len(prs.slides)} slides at {width_px}x{height_px}")
        
        for i, slide in enumerate(prs.slides):
            output_path = os.path.join(self.output_dir, f"slide_{i:03d}.png")
            
            # Create image with light gray background
            img = Image.new('RGB', (width_px, height_px), '#f0f0f0')
            draw = ImageDraw.Draw(img)
            
            # Try to load a default font
            try:
                font = ImageFont.truetype("arial.ttf", 24)
                font_small = ImageFont.truetype("arial.ttf", 18)
            except:
                font = ImageFont.load_default()
                font_small = ImageFont.load_default()
            
            # Extract and draw title
            title_y = 50
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    try:
                        # Get position
                        left = int(shape.left / 914400 * self.dpi) if hasattr(shape, "left") else 50
                        top = int(shape.top / 914400 * self.dpi) if hasattr(shape, "top") else 50
                        width = int(shape.width / 914400 * self.dpi) if hasattr(shape, "width") else width_px - 100
                        height = int(shape.height / 914400 * self.dpi) if hasattr(shape, "height") else 100
                        
                        # Draw text
                        text = shape.text[:500]  # Limit text length
                        draw.text((left, top), text, fill='black', font=font)
                    except:
                        pass
            
            # Draw slide number
            draw.text((width_px - 80, height_px - 40), f"Slide {i+1}", fill='gray', font=font_small)
            
            img.save(output_path, "PNG")
            image_paths.append(output_path)
            print(f"  Exported slide {i+1}/{len(prs.slides)}")
        
        return image_paths


class VideoGenerator:
    def __init__(self, output_dir: str, fps: int = 30, slide_duration: float = 5.0):
        self.output_dir = output_dir
        self.fps = fps
        self.slide_duration = slide_duration

    def create_video_from_slides_and_audio(
        self, 
        slide_images: list[str], 
        audio_files: list[Optional[str]], 
        output_path: str
    ) -> str:
        import subprocess
        import glob
        
        temp_dir = self.output_dir
        
        # Build list of input files with durations
        inputs = []
        for i, slide_img in enumerate(slide_images):
            if not slide_img or not os.path.exists(slide_img):
                continue
            
            # Calculate duration
            duration = self.slide_duration
            audio_file = audio_files[i] if i < len(audio_files) else None
            if audio_file and os.path.exists(audio_file):
                from moviepy import AudioFileClip
                audio_clip = AudioFileClip(audio_file)
                duration = max(duration, audio_clip.duration)
                audio_clip.close()
            
            inputs.append((slide_img, duration))
        
        # Create video using image sequence with loop
        video_no_audio = os.path.join(temp_dir, "video_no_audio.mp4")
        
        # Use a simple approach: create video from each image with its duration
        parts = []
        for img_path, duration in inputs:
            part_file = os.path.join(temp_dir, f"part_{len(parts)}.mp4")
            cmd = [
                'ffmpeg', '-y',
                '-loop', '1',
                '-i', img_path,
                '-c:v', 'libx264',
                '-preset', 'medium',
                '-crf', '23',
                '-t', str(duration),
                '-vf', f"scale=1280:720:force_original_aspect_ratio=decrease,pad=1280:720:(ow-iw)/2:(oh-ih)/2",
                '-pix_fmt', 'yuv420p',
                '-r', '30',
                '-fps_mode', 'cfr',
                part_file
            ]
            result = subprocess.run(cmd, capture_output=True, text=True)
            if result.returncode != 0:
                print(f"FFmpeg part error: {result.stderr[:500]}")
                continue
            parts.append(part_file)
        
        # Concatenate all parts
        if len(parts) == 1:
            final_video = parts[0]
        else:
            # Write concat list
            concat_file = os.path.join(temp_dir, "concat.txt")
            with open(concat_file, 'w') as f:
                for p in parts:
                    f.write(f"file '{p}'\n")
            
            final_video = os.path.join(temp_dir, "video_concat.mp4")
            cmd = [
                'ffmpeg', '-y',
                '-f', 'concat',
                '-safe', '0',
                '-i', concat_file,
                '-c', 'copy',
                final_video
            ]
            result = subprocess.run(cmd, capture_output=True, text=True)
            if result.returncode != 0:
                print(f"FFmpeg concat error: {result.stderr[:500]}")
                final_video = parts[0]
        
        # Check if we have any audio
        has_audio = any(a and os.path.exists(a) for a in audio_files if a)
        
        if has_audio:
            # Combine all audio files into one
            audio_concat = os.path.join(temp_dir, "audio_concat.txt")
            audio_lines = []
            for a in audio_files:
                if a and os.path.exists(a):
                    audio_lines.append(f"file '{a}'")
            
            with open(audio_concat, 'w') as f:
                f.write("\n".join(audio_lines))
            
            combined_audio = os.path.join(temp_dir, "combined_audio.m4a")
            
            cmd_audio = [
                'ffmpeg', '-y',
                '-f', 'concat',
                '-safe', '0',
                '-i', audio_concat,
                '-c:a', 'aac',
                '-b:a', '192k',
                combined_audio
            ]
            result = subprocess.run(cmd_audio, capture_output=True, text=True)
            
            # Now combine video and audio
            cmd_final = [
                'ffmpeg', '-y',
                '-i', final_video,
                '-i', combined_audio,
                '-c:v', 'copy',
                '-c:a', 'aac',
                '-shortest',
                output_path
            ]
            result = subprocess.run(cmd_final, capture_output=True, text=True)
            if result.returncode != 0:
                print(f"FFmpeg final error: {result.stderr[:500]}")
                # Fall back to just video
                import shutil
                shutil.copy(final_video, output_path)
        else:
            # Just copy video
            import shutil
            shutil.copy(final_video, output_path)
        
        return output_path


class PPT2Video:
    def __init__(
        self,
        pptx_path: str,
        output_path: Optional[str] = None,
        lang: str = "en",
        slide_duration: float = 5.0,
        fps: int = 30,
        dpi: int = 150
    ):
        self.pptx_path = pptx_path
        self.output_path = output_path
        self.lang = lang
        self.slide_duration = slide_duration
        self.fps = fps
        self.dpi = dpi
        
        # Create temporary directory for intermediate files
        self.temp_dir = tempfile.mkdtemp()
        self.images_dir = os.path.join(self.temp_dir, "images")
        self.audio_dir = os.path.join(self.temp_dir, "audio")
        
        os.makedirs(self.images_dir, exist_ok=True)
        os.makedirs(self.audio_dir, exist_ok=True)

    def run(self) -> str:
        click.echo("Reading PowerPoint file...")
        reader = PowerPointReader(self.pptx_path)
        slides_data = reader.extract_slides_and_notes()
        
        click.echo(f"Found {len(slides_data)} slides")
        
        # Convert slides to images
        click.echo("Exporting slides to images...")
        slide_exporter = SlideToImage(self.pptx_path, self.images_dir, self.dpi)
        slide_images = slide_exporter.export_slides()
        
        # Generate audio for each slide
        click.echo("Generating voiceover audio...")
        tts = TextToSpeech(self.audio_dir, self.lang)
        audio_files = []
        
        for i, slide_data in enumerate(slides_data):
            notes = slide_data["notes"]
            if notes:
                audio_path = tts.text_to_audio(notes, f"slide_{i:03d}.mp3")
                audio_files.append(audio_path)
            else:
                audio_files.append(None)
        
        # Create video
        click.echo("Creating video...")
        video_generator = VideoGenerator(
            self.temp_dir,
            fps=self.fps,
            slide_duration=self.slide_duration
        )
        
        output_path = self.output_path or os.path.join(
            os.path.dirname(self.pptx_path),
            os.path.splitext(os.path.basename(self.pptx_path))[0] + ".mp4"
        )
        
        video_generator.create_video_from_slides_and_audio(
            slide_images,
            audio_files,
            output_path
        )
        
        # Cleanup temp files
        shutil.rmtree(self.temp_dir)
        
        click.echo(f"Video saved to: {output_path}")
        return output_path


@click.command()
@click.argument("pptx_file", type=click.Path(exists=True))
@click.option("-o", "--output", "output_path", type=click.Path(), help="Output video file path")
@click.option("-l", "--lang", default="en", help="Language for text-to-speech (e.g., en, es, fr)")
@click.option("-d", "--duration", default=5.0, type=float, help="Minimum duration per slide in seconds")
@click.option("--fps", default=30, type=int, help="Frames per second for output video")
@click.option("--dpi", default=150, type=int, help=" DPI for slide images")
def cli(pptx_file: str, output_path: Optional[str], lang: str, duration: float, fps: int, dpi: int):
    """Convert a PowerPoint presentation to video with voiceover."""
    converter = PPT2Video(
        pptx_path=pptx_file,
        output_path=output_path,
        lang=lang,
        slide_duration=duration,
        fps=fps,
        dpi=dpi
    )
    converter.run()


if __name__ == "__main__":
    cli()
